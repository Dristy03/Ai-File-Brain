from __future__ import annotations

import asyncio
import logging

from ai_file_brain.core.models import ExtractionResult

logger = logging.getLogger(__name__)


class PptxExtractor:
    async def extract(self, file_path: str) -> ExtractionResult:
        text = await asyncio.to_thread(self._read_sync, file_path)
        return ExtractionResult(text=text, source="native")

    @staticmethod
    def _read_sync(file_path: str) -> str:
        try:
            from pptx import Presentation
        except ImportError as ex:
            logger.warning("python-pptx not available; cannot read %s: %s", file_path, ex)
            return ""

        try:
            prs = Presentation(file_path)
        except FileNotFoundError:
            return ""
        except Exception as ex:
            # Password-protected, malformed, or a legacy .ppt misnamed as .pptx.
            logger.warning("Failed to open PPTX %s: %s", file_path, ex)
            return ""

        parts: list[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                # Text boxes, titles, body placeholders.
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = "".join(run.text for run in para.runs)
                        if line.strip():
                            parts.append(line)
                # Tables on the slide.
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [cell.text for cell in row.cells if cell.text.strip()]
                        if cells:
                            parts.append("\t".join(cells))
            # Speaker notes.
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame
                if notes is not None and notes.text.strip():
                    parts.append(notes.text)

        return "\n".join(parts)
