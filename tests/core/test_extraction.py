from __future__ import annotations

from pathlib import Path
from typing import Iterable

import pytest

from ai_file_brain.core.extraction import (
    UnsupportedFileTypeError,
    get_extractor,
    is_supported,
)
from ai_file_brain.core.models import ExtractionResult


def test_is_supported():
    assert is_supported("a.TXT")
    assert is_supported("b.pdf")
    assert is_supported("c.PNG")
    assert is_supported("d.jpg")
    assert is_supported("e.jpeg")
    assert is_supported("f.tif")
    assert is_supported("g.tiff")
    assert is_supported("h.bmp")
    assert is_supported("i.webp")
    assert is_supported("j.docx")
    # Code / config — share PlainTextExtractor.
    assert is_supported("k.py")
    assert is_supported("l.JSON")
    assert is_supported("m.yaml")
    assert is_supported("n.toml")
    assert is_supported("o.md")
    assert is_supported("r.pptx")
    assert is_supported("s.XLSX")
    assert is_supported("t.ppt")  # legacy binary PowerPoint
    assert is_supported("u.xls")  # legacy binary Excel
    assert is_supported("v.DOC")  # legacy binary Word
    assert not is_supported("p.exe")


def test_unsupported_raises():
    with pytest.raises(UnsupportedFileTypeError):
        get_extractor("/whatever/file.exe")


@pytest.mark.asyncio
async def test_plain_text_extractor_reads_file(tmp_path: Path):
    file = tmp_path / "hello.txt"
    file.write_text("greetings, earth", encoding="utf-8")
    extractor = get_extractor(str(file))
    result = await extractor.extract(str(file))
    assert isinstance(result, ExtractionResult)
    assert result.text == "greetings, earth"
    assert result.source == "native"


@pytest.mark.asyncio
async def test_pdf_extractor_returns_empty_for_missing_file(tmp_path: Path):
    fake = tmp_path / "missing.pdf"
    extractor = get_extractor(str(fake))
    result = await extractor.extract(str(fake))
    assert result.text == ""
    assert result.source == "native"


# --- code / config files (share PlainTextExtractor) ---


@pytest.mark.asyncio
async def test_code_file_routes_to_plain_text(tmp_path: Path):
    py_file = tmp_path / "snippet.py"
    py_file.write_text("def hello():\n    return 'world'\n", encoding="utf-8")
    extractor = get_extractor(str(py_file))
    result = await extractor.extract(str(py_file))
    assert result.source == "native"
    assert "def hello" in result.text


@pytest.mark.asyncio
async def test_json_file_routes_to_plain_text(tmp_path: Path):
    j = tmp_path / "config.json"
    j.write_text('{"name": "ai-file-brain", "version": "0.2"}', encoding="utf-8")
    extractor = get_extractor(str(j))
    result = await extractor.extract(str(j))
    assert result.source == "native"
    assert "ai-file-brain" in result.text


# --- docx ---


def _make_docx(path: Path, paragraphs: list[str]) -> None:
    from docx import Document

    doc = Document()
    for para in paragraphs:
        doc.add_paragraph(para)
    doc.save(str(path))


def _make_docx_with_table_and_header(
    path: Path, body: list[str], table: list[list[str]], header: str
) -> None:
    from docx import Document

    doc = Document()
    for para in body:
        doc.add_paragraph(para)
    if table:
        rows = len(table)
        cols = len(table[0])
        t = doc.add_table(rows=rows, cols=cols)
        for r, row_vals in enumerate(table):
            for c, val in enumerate(row_vals):
                t.cell(r, c).text = val
    if header:
        section = doc.sections[0]
        section.header.paragraphs[0].text = header
    doc.save(str(path))


@pytest.mark.asyncio
async def test_docx_extractor_reads_paragraphs(tmp_path: Path):
    path = tmp_path / "simple.docx"
    _make_docx(path, ["First paragraph.", "Second paragraph with detail."])
    extractor = get_extractor(str(path))
    result = await extractor.extract(str(path))
    assert result.source == "native"
    assert "First paragraph." in result.text
    assert "Second paragraph with detail." in result.text


@pytest.mark.asyncio
async def test_docx_extractor_reads_tables_and_header(tmp_path: Path):
    path = tmp_path / "rich.docx"
    _make_docx_with_table_and_header(
        path,
        body=["Body text here."],
        table=[["Header A", "Header B"], ["Cell A1", "Cell B1"]],
        header="Confidential — internal only",
    )
    extractor = get_extractor(str(path))
    result = await extractor.extract(str(path))
    assert result.source == "native"
    assert "Body text here." in result.text
    assert "Header A" in result.text
    assert "Cell B1" in result.text
    assert "Confidential" in result.text


@pytest.mark.asyncio
async def test_docx_extractor_handles_corrupt_file(tmp_path: Path):
    path = tmp_path / "broken.docx"
    path.write_bytes(b"this is not a real docx")
    extractor = get_extractor(str(path))
    result = await extractor.extract(str(path))
    assert result.text == ""
    assert result.source == "native"


@pytest.mark.asyncio
async def test_docx_extractor_handles_missing_file(tmp_path: Path):
    fake = tmp_path / "ghost.docx"
    extractor = get_extractor(str(fake))
    result = await extractor.extract(str(fake))
    assert result.text == ""
    assert result.source == "native"


# --- pptx ---


def _make_pptx(path: Path, title: str, body: str, notes: str = "") -> None:
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # title + content
    slide.shapes.title.text = title
    slide.placeholders[1].text = body
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    prs.save(str(path))


@pytest.mark.asyncio
async def test_pptx_extractor_reads_slides_and_notes(tmp_path: Path):
    path = tmp_path / "deck.pptx"
    _make_pptx(path, "Q3 Roadmap", "Ship the new indexer", notes="Mention the GPU work")
    extractor = get_extractor(str(path))
    result = await extractor.extract(str(path))
    assert result.source == "native"
    assert "Q3 Roadmap" in result.text
    assert "Ship the new indexer" in result.text
    assert "Mention the GPU work" in result.text


@pytest.mark.asyncio
async def test_pptx_extractor_handles_corrupt_file(tmp_path: Path):
    path = tmp_path / "broken.pptx"
    path.write_bytes(b"not a real pptx")
    extractor = get_extractor(str(path))
    result = await extractor.extract(str(path))
    assert result.text == ""
    assert result.source == "native"


# --- xlsx ---


def _make_xlsx(path: Path, sheets: dict[str, list[list[object]]]) -> None:
    from openpyxl import Workbook

    wb = Workbook()
    wb.remove(wb.active)  # drop the default empty sheet
    for name, rows in sheets.items():
        ws = wb.create_sheet(title=name)
        for row in rows:
            ws.append(row)
    wb.save(str(path))


@pytest.mark.asyncio
async def test_xlsx_extractor_reads_cells_and_sheet_titles(tmp_path: Path):
    path = tmp_path / "book.xlsx"
    _make_xlsx(
        path,
        {
            "May": [["Name", "Hours"], ["Asha", 40], ["Ben", 38]],
            "Notes": [["Reconcile with HR"]],
        },
    )
    extractor = get_extractor(str(path))
    result = await extractor.extract(str(path))
    assert result.source == "native"
    assert "May" in result.text
    assert "Asha" in result.text
    assert "40" in result.text
    assert "Reconcile with HR" in result.text


@pytest.mark.asyncio
async def test_xlsx_extractor_handles_corrupt_file(tmp_path: Path):
    path = tmp_path / "broken.xlsx"
    path.write_bytes(b"not a real xlsx")
    extractor = get_extractor(str(path))
    result = await extractor.extract(str(path))
    assert result.text == ""
    assert result.source == "native"


# --- legacy .xls (xlrd) ---


def _make_xls(path: Path, sheet_name: str, rows: list[list[object]]) -> None:
    import xlwt

    wb = xlwt.Workbook()
    ws = wb.add_sheet(sheet_name)
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            ws.write(r, c, val)
    wb.save(str(path))


@pytest.mark.asyncio
async def test_xls_extractor_reads_cells(tmp_path: Path):
    path = tmp_path / "legacy.xls"
    _make_xls(path, "Sheet1", [["Name", "Hours"], ["Asha", 40], ["Ben", 38]])
    extractor = get_extractor(str(path))
    result = await extractor.extract(str(path))
    assert result.source == "native"
    assert "Sheet1" in result.text
    assert "Asha" in result.text
    assert "40" in result.text


@pytest.mark.asyncio
async def test_xls_extractor_handles_corrupt_file(tmp_path: Path):
    path = tmp_path / "broken.xls"
    path.write_bytes(b"not a real xls")
    extractor = get_extractor(str(path))
    result = await extractor.extract(str(path))
    assert result.text == ""
    assert result.source == "native"


# --- legacy .ppt (olefile + atom parsing) ---


def _ppt_record(ver_inst: int, rec_type: int, payload: bytes) -> bytes:
    import struct

    return struct.pack("<HHI", ver_inst, rec_type, len(payload)) + payload


def test_extract_ppt_text_decodes_text_atoms():
    from ai_file_brain.core.extraction.ppt import extract_ppt_text

    # A TextBytesAtom (0x0FA8, Latin-1) and a TextCharsAtom (0x0FA0, UTF-16LE),
    # both nested inside a container record (recVer nibble == 0xF).
    bytes_atom = _ppt_record(0x0000, 0x0FA8, b"Hello legacy slide")
    chars_atom = _ppt_record(0x0000, 0x0FA0, "Second line".encode("utf-16-le"))
    container = _ppt_record(0x000F, 0x0FF0, bytes_atom + chars_atom)

    text = extract_ppt_text(container)
    assert "Hello legacy slide" in text
    assert "Second line" in text


def test_extract_ppt_text_handles_garbage():
    from ai_file_brain.core.extraction.ppt import extract_ppt_text

    assert extract_ppt_text(b"\x00\x01\x02not really records") == ""


@pytest.mark.asyncio
async def test_ppt_extractor_handles_non_ole_file(tmp_path: Path):
    path = tmp_path / "broken.ppt"
    path.write_bytes(b"not an OLE2 file")
    extractor = get_extractor(str(path))
    result = await extractor.extract(str(path))
    assert result.text == ""
    assert result.source == "native"


@pytest.mark.asyncio
async def test_ppt_extractor_handles_missing_file(tmp_path: Path):
    fake = tmp_path / "ghost.ppt"
    extractor = get_extractor(str(fake))
    result = await extractor.extract(str(fake))
    assert result.text == ""
    assert result.source == "native"


# --- legacy .doc (olefile + piece-table parsing) ---


def _build_doc_streams(pieces: list[tuple[str, bool]]) -> tuple[bytes, bytes]:
    """Hand-build minimal WordDocument + table streams for a list of
    (text, compressed) pieces, exercising the real piece-table parser."""
    import struct

    word = bytearray(0x2000)
    struct.pack_into("<H", word, 0x0000, 0xA5EC)  # wIdent
    struct.pack_into("<H", word, 0x000A, 0x0000)  # flags -> 0Table

    cps = [0]
    pcds = bytearray()
    text_off = 0x400
    for text, compressed in pieces:
        cch = len(text)
        cps.append(cps[-1] + cch)
        if compressed:
            word[text_off : text_off + cch] = text.encode("cp1252")
            fc_raw = 0x40000000 | (text_off * 2)
            text_off += cch + 4
        else:
            raw = text.encode("utf-16-le")
            word[text_off : text_off + len(raw)] = raw
            fc_raw = text_off
            text_off += len(raw) + 4
        pcds += struct.pack("<HIH", 0, fc_raw, 0)  # flags, fc, prm

    plcfpcd = b"".join(struct.pack("<I", cp) for cp in cps) + bytes(pcds)
    clx = bytes([0x02]) + struct.pack("<I", len(plcfpcd)) + plcfpcd
    struct.pack_into("<I", word, 0x01A2, 0)  # fcClx
    struct.pack_into("<I", word, 0x01A6, len(clx))  # lcbClx
    return bytes(word), clx


def test_extract_doc_text_stitches_compressed_and_unicode_pieces():
    from ai_file_brain.core.extraction.doc import extract_doc_text

    word, table = _build_doc_streams([("Hello DOC ", True), ("wide piece", False)])
    text = extract_doc_text(word, table)
    assert "Hello DOC" in text
    assert "wide piece" in text


def test_extract_doc_text_rejects_non_word_stream():
    from ai_file_brain.core.extraction.doc import extract_doc_text

    assert extract_doc_text(b"\x00" * 500, b"") == ""


def test_doc_table_stream_name_follows_flag():
    import struct

    from ai_file_brain.core.extraction.doc import table_stream_name

    word = bytearray(16)
    struct.pack_into("<H", word, 0x0A, 0x0200)
    assert table_stream_name(bytes(word)) == "1Table"
    struct.pack_into("<H", word, 0x0A, 0x0000)
    assert table_stream_name(bytes(word)) == "0Table"


@pytest.mark.asyncio
async def test_doc_extractor_handles_non_ole_file(tmp_path: Path):
    path = tmp_path / "broken.doc"
    path.write_bytes(b"not an OLE2 file")
    extractor = get_extractor(str(path))
    result = await extractor.extract(str(path))
    assert result.text == ""
    assert result.source == "native"


@pytest.mark.asyncio
async def test_doc_extractor_handles_missing_file(tmp_path: Path):
    fake = tmp_path / "ghost.doc"
    extractor = get_extractor(str(fake))
    result = await extractor.extract(str(fake))
    assert result.text == ""
    assert result.source == "native"


# --- helpers for OCR tests ---


def _make_text_image(text: str, size: tuple[int, int] = (800, 200)):
    from PIL import Image, ImageDraw, ImageFont

    img = Image.new("RGB", size, "white")
    draw = ImageDraw.Draw(img)
    font = None
    for candidate in ("arial.ttf", "DejaVuSans.ttf", "FreeSans.ttf"):
        try:
            font = ImageFont.truetype(candidate, 64)
            break
        except OSError:
            continue
    if font is None:
        font = ImageFont.load_default()
    draw.text((40, 50), text, fill="black", font=font)
    return img


def _make_native_pdf(path: Path, lines: Iterable[str]) -> None:
    import pymupdf

    doc = pymupdf.open()
    page = doc.new_page()
    y = 72.0
    for line in lines:
        page.insert_text((72, y), line, fontsize=18)
        y += 24
    doc.save(str(path))
    doc.close()


def _make_image_pdf(path: Path, text: str) -> None:
    import pymupdf

    img = _make_text_image(text)
    img_path = path.with_suffix(".helper.png")
    img.save(img_path)
    try:
        doc = pymupdf.open()
        page = doc.new_page(width=612, height=792)
        rect = pymupdf.Rect(36, 36, 576, 756)
        page.insert_image(rect, filename=str(img_path))
        doc.save(str(path))
        doc.close()
    finally:
        try:
            img_path.unlink()
        except OSError:
            pass


def _make_mixed_pdf(path: Path, native_line: str, image_text: str) -> None:
    import pymupdf

    img = _make_text_image(image_text)
    img_path = path.with_suffix(".helper.png")
    img.save(img_path)
    try:
        doc = pymupdf.open()
        # Page 1: native text
        page1 = doc.new_page()
        page1.insert_text((72, 72), native_line, fontsize=18)
        # Page 2: image only
        page2 = doc.new_page(width=612, height=792)
        rect = pymupdf.Rect(36, 36, 576, 756)
        page2.insert_image(rect, filename=str(img_path))
        doc.save(str(path))
        doc.close()
    finally:
        try:
            img_path.unlink()
        except OSError:
            pass


def _ocr_match(haystack: str, needle: str) -> bool:
    """Loose match: OCR may emit slightly different casing or extra whitespace."""
    return needle.lower().replace(" ", "") in haystack.lower().replace(" ", "")


# --- image extractor ---


@pytest.mark.slow
@pytest.mark.asyncio
async def test_image_extractor_reads_text_from_png(tmp_path: Path):
    img_path = tmp_path / "hello.png"
    _make_text_image("HELLO WORLD").save(img_path)
    extractor = get_extractor(str(img_path))
    result = await extractor.extract(str(img_path))
    assert result.source == "ocr"
    assert _ocr_match(result.text, "HELLO") or _ocr_match(result.text, "WORLD")


@pytest.mark.slow
@pytest.mark.asyncio
async def test_image_extractor_returns_empty_for_blank_image(tmp_path: Path):
    from PIL import Image

    img_path = tmp_path / "blank.png"
    Image.new("RGB", (300, 300), "white").save(img_path)
    extractor = get_extractor(str(img_path))
    result = await extractor.extract(str(img_path))
    assert result.text == ""
    assert result.source == "ocr"


@pytest.mark.asyncio
async def test_image_extractor_handles_corrupt_file(tmp_path: Path):
    bogus = tmp_path / "broken.png"
    bogus.write_bytes(b"this is not a real image")
    extractor = get_extractor(str(bogus))
    result = await extractor.extract(str(bogus))
    assert result.text == ""
    assert result.source == "ocr"


@pytest.mark.asyncio
async def test_image_extractor_missing_file(tmp_path: Path):
    fake = tmp_path / "ghost.png"
    extractor = get_extractor(str(fake))
    result = await extractor.extract(str(fake))
    assert result.text == ""
    assert result.source == "ocr"


# --- pdf extractor ---


@pytest.mark.asyncio
async def test_pdf_extractor_native_fast_path(tmp_path: Path):
    pdf_path = tmp_path / "native.pdf"
    _make_native_pdf(
        pdf_path,
        [
            "The quick brown fox jumps over the lazy dog.",
            "Pack my box with five dozen liquor jugs.",
        ],
    )
    extractor = get_extractor(str(pdf_path))
    result = await extractor.extract(str(pdf_path))
    assert result.source == "native"
    assert "quick brown fox" in result.text


@pytest.mark.slow
@pytest.mark.asyncio
async def test_pdf_extractor_image_only_pdf_uses_ocr(tmp_path: Path):
    pdf_path = tmp_path / "scan.pdf"
    _make_image_pdf(pdf_path, "HELLO WORLD")
    extractor = get_extractor(str(pdf_path))
    result = await extractor.extract(str(pdf_path))
    assert result.source == "ocr"
    assert _ocr_match(result.text, "HELLO") or _ocr_match(result.text, "WORLD")


@pytest.mark.slow
@pytest.mark.asyncio
async def test_pdf_extractor_mixed_pages_returns_mixed_source(tmp_path: Path):
    pdf_path = tmp_path / "mixed.pdf"
    native_line = "Native page content goes here for the test."
    _make_mixed_pdf(pdf_path, native_line, "HELLO WORLD")
    extractor = get_extractor(str(pdf_path))
    result = await extractor.extract(str(pdf_path))
    # The whole-doc native chars from page 1 alone may exceed the threshold.
    # In that case, the fast path is taken and the image page contributes nothing.
    # Otherwise, fallback runs and source is "mixed". Both behaviours are valid;
    # the assertion captures the one we configured for in the spec.
    assert "Native page content" in result.text
    if result.source == "mixed":
        assert _ocr_match(result.text, "HELLO") or _ocr_match(result.text, "WORLD")


@pytest.mark.slow
@pytest.mark.asyncio
async def test_pdf_extractor_mixed_pages_fallback_when_threshold_high(
    tmp_path: Path, monkeypatch
):
    """Force the fallback path by raising the doc-level threshold above the page-1 length."""
    monkeypatch.setenv("AFB_PDF_OCR_MIN_NATIVE_CHARS", "5000")
    pdf_path = tmp_path / "mixed_forced.pdf"
    _make_mixed_pdf(pdf_path, "Native page content.", "HELLO WORLD")
    extractor = get_extractor(str(pdf_path))
    result = await extractor.extract(str(pdf_path))
    assert result.source == "mixed"
    assert "Native page content" in result.text
    assert _ocr_match(result.text, "HELLO") or _ocr_match(result.text, "WORLD")


# --- ocr disabled escape hatch ---


@pytest.mark.asyncio
async def test_ocr_disabled_image_returns_empty(tmp_path: Path, monkeypatch):
    monkeypatch.setenv("AFB_OCR_ENABLED", "false")
    img_path = tmp_path / "skip.png"
    _make_text_image("HELLO").save(img_path)
    extractor = get_extractor(str(img_path))
    result = await extractor.extract(str(img_path))
    assert result.text == ""
    assert result.source == "ocr"


@pytest.mark.asyncio
async def test_ocr_disabled_scanned_pdf_takes_native_fast_path(
    tmp_path: Path, monkeypatch
):
    monkeypatch.setenv("AFB_OCR_ENABLED", "false")
    pdf_path = tmp_path / "scan_disabled.pdf"
    _make_image_pdf(pdf_path, "HELLO WORLD")
    extractor = get_extractor(str(pdf_path))
    result = await extractor.extract(str(pdf_path))
    assert result.source == "native"
    # Image-only PDF + OCR disabled → no extracted text.
    assert result.text.strip() == ""
