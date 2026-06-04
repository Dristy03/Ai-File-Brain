"""Generate one sample file per supported type into ./sample-files.

Point the app's watch folder at the generated directory to exercise every
extractor / indexing tier. Re-run any time; it overwrites in place.

Gitignored (see .gitignore: sample-files/). Run with the project venv:
    .venv\\Scripts\\python.exe scripts\\gen_sample_files.py
"""

from __future__ import annotations

import os
import zipfile
from pathlib import Path

OUT = Path(__file__).resolve().parent.parent / "sample-files"

# A common phrase salted into every file so semantic / substring search has
# something predictable to match across types.
MARKER = "AI File Brain sample fixture. The secret marker phrase is bluefin-koala-2026."


def _plain_text_files() -> None:
    """Content tier: plain-text + source-code + config formats, all UTF-8."""
    samples: dict[str, str] = {
        "notes.txt": f"{MARKER}\nThis is a plain text note about quarterly planning.\n",
        "readme.md": f"# Sample Markdown\n\n{MARKER}\n\n- bullet one\n- bullet two\n",
        "manual.rst": f"Sample reStructuredText\n=======================\n\n{MARKER}\n",
        "script.py": f'"""{MARKER}"""\n\n\ndef greet(name):\n    return f"hello {{name}}"\n',
        "app.js": f"// {MARKER}\nfunction greet(name) {{ return `hello ${{name}}`; }}\n",
        "types.ts": f"// {MARKER}\nexport const greet = (name: string): string => `hi ${{name}}`;\n",
        "view.tsx": f"// {MARKER}\nexport const View = () => <div>hello</div>;\n",
        "widget.jsx": f"// {MARKER}\nexport const Widget = () => <span>hi</span>;\n",
        "Main.java": f"// {MARKER}\nclass Main {{ public static void main(String[] a) {{}} }}\n",
        "Program.cs": f"// {MARKER}\nclass Program {{ static void Main() {{}} }}\n",
        "main.go": f"// {MARKER}\npackage main\n\nfunc main() {{}}\n",
        "lib.rs": f"// {MARKER}\npub fn greet() -> &'static str {{ \"hi\" }}\n",
        "app.rb": f"# {MARKER}\ndef greet(name) = \"hello #{{name}}\"\n",
        "index.php": f"<?php // {MARKER}\nfunction greet($n) {{ return \"hi $n\"; }}\n",
        "main.c": f"/* {MARKER} */\nint main(void) {{ return 0; }}\n",
        "main.cpp": f"// {MARKER}\nint main() {{ return 0; }}\n",
        "util.cc": f"// {MARKER}\nint util() {{ return 1; }}\n",
        "header.h": f"/* {MARKER} */\n#pragma once\nint api(void);\n",
        "header.hpp": f"// {MARKER}\n#pragma once\nclass Api {{}};\n",
        "deploy.sh": f"#!/bin/sh\n# {MARKER}\necho hello\n",
        "setup.bash": f"#!/usr/bin/env bash\n# {MARKER}\necho hi\n",
        "task.ps1": f"# {MARKER}\nWrite-Output 'hello'\n",
        "query.sql": f"-- {MARKER}\nSELECT 1 AS answer;\n",
        # config / data text formats (indexed via plain-text extractor)
        "config.yml": f"# {MARKER}\nname: sample\nenabled: true\n",
        "config.yaml": f"# {MARKER}\nname: sample\n",
        "pyproject_demo.toml": f'# {MARKER}\n[tool.demo]\nname = "sample"\n',
        "data.json": f'{{"marker": "{MARKER}", "items": [1, 2, 3]}}\n',
        "settings.ini": f"; {MARKER}\n[section]\nkey = value\n",
        "app.cfg": f"# {MARKER}\n[app]\ndebug = false\n",
        "dotenv.env": f"# {MARKER}\nAPP_NAME=sample\n",
    }
    for name, body in samples.items():
        (OUT / name).write_text(body, encoding="utf-8")


def _pdf() -> None:
    import fitz  # PyMuPDF

    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((72, 72), f"Sample PDF document.\n{MARKER}")
    doc.save(str(OUT / "document.pdf"))
    doc.close()


def _docx() -> None:
    from docx import Document

    doc = Document()
    doc.add_heading("Sample Word Document", level=1)
    doc.add_paragraph(MARKER)
    doc.add_paragraph("A second paragraph about project status.")
    doc.save(str(OUT / "report.docx"))


def _pptx() -> None:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Sample Slide Deck"
    box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
    box.text_frame.text = MARKER
    prs.save(str(OUT / "deck.pptx"))


def _xlsx() -> None:
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    ws.append(["marker", "value"])
    ws.append([MARKER, 42])
    ws.append(["quarter", "Q2-2026"])
    wb.save(str(OUT / "spreadsheet.xlsx"))


def _images() -> None:
    from PIL import Image, ImageDraw

    # Name-only tier (OCR off by default for images), so pixel content is just a
    # placeholder; the filename is what gets indexed. One per supported ext.
    exts = ["png", "jpg", "jpeg", "tif", "tiff", "bmp", "webp", "gif"]
    base = Image.new("RGB", (240, 80), color=(30, 90, 160))
    draw = ImageDraw.Draw(base)
    draw.text((10, 30), "sample image", fill=(255, 255, 255))
    for ext in exts:
        img = base.convert("P") if ext == "gif" else base
        img.save(str(OUT / f"screenshot.{ext}"))


def _zip() -> None:
    with zipfile.ZipFile(OUT / "archive.zip", "w", zipfile.ZIP_DEFLATED) as zf:
        zf.writestr("inside/readme.txt", f"zipped content\n{MARKER}\n")


def _mp4_placeholder() -> None:
    # Name-only tier — content is never read, so a tiny placeholder with a valid
    # ftyp box header is plenty to represent the type on disk.
    ftyp = bytes([
        0x00, 0x00, 0x00, 0x18, 0x66, 0x74, 0x79, 0x70,  # box size + 'ftyp'
        0x69, 0x73, 0x6F, 0x6D, 0x00, 0x00, 0x02, 0x00,  # major brand 'isom'
        0x69, 0x73, 0x6F, 0x6D, 0x69, 0x73, 0x6F, 0x32,  # compatible brands
    ])
    (OUT / "clip.mp4").write_bytes(ftyp)


def _legacy_note() -> None:
    # .doc / .ppt / .xls are legacy OLE binary formats that can't be authored
    # with the available Python libs (and LibreOffice isn't installed). Drop a
    # note so it's clear they're intentionally absent, not forgotten.
    (OUT / "_LEGACY_FORMATS_README.txt").write_text(
        "Legacy binary Office formats (.doc, .ppt, .xls) are NOT generated here:\n"
        "they require real OLE2 files, which can't be produced by the installed\n"
        "Python libraries. To test those extractors, open report.docx / deck.pptx\n"
        "/ spreadsheet.xlsx in Office (or LibreOffice) and 'Save As' the 97-2003\n"
        "formats into this folder as report.doc / deck.ppt / spreadsheet.xls.\n",
        encoding="utf-8",
    )


def main() -> None:
    OUT.mkdir(parents=True, exist_ok=True)
    _plain_text_files()
    _pdf()
    _docx()
    _pptx()
    _xlsx()
    _images()
    _zip()
    _mp4_placeholder()
    _legacy_note()
    count = sum(1 for _ in OUT.iterdir() if _.is_file())
    print(f"Generated {count} files in {OUT}")


if __name__ == "__main__":
    main()
